from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MEPI Methodology Overview"
    subtitle.text = "Energy Poverty Mapping Initiative"

    # Content slides
    methodologies = [
        "1. Data Collection: Gathering data on energy usage and poverty metrics.",
        "2. Analysis: Assessing the relationship between energy access and poverty levels.",
        "3. Recommendations: Proposing initiatives based on the analysis."
    ]
    
    for methodology in methodologies:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Methodology Step"
        content.text = methodology
    
    # Save the presentation
    prs.save('MEPI_Methodology_Presentation.pptx')

if __name__ == "__main__":
    create_presentation()